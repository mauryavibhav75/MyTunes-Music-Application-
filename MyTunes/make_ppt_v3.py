"""
MyTunes – Final Presentation (v3)
Slide order:
  1  Title
  2  Introduction          ← NEW
  3  Problem Statement     ← NEW
  4  Objectives            ← NEW
  5  Scope of the Project  ← NEW
  6  Technologies Used     ← NEW
  7  Project Overview      (existing)
  8–12 Screenshot slides   (existing)
  13 System Architecture   (existing)
  14 Outcomes / Results    ← NEW
  15 Future Work           ← NEW
  16 Thank You             (existing)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Colours ───────────────────────────────────────────────────────────────────
BG_DARK  = RGBColor(0x12, 0x12, 0x12)
BG_CARD  = RGBColor(0x1E, 0x1E, 0x2E)
ACCENT   = RGBColor(0x1D, 0xB9, 0x54)   # green
ACCENT2  = RGBColor(0x9B, 0x59, 0xB6)   # purple
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY     = RGBColor(0xB3, 0xB3, 0xB3)
YELLOW   = RGBColor(0xF3, 0xC6, 0x23)
TEAL     = RGBColor(0x1A, 0xBC, 0x9C)
RED      = RGBColor(0xE5, 0x3E, 0x3E)
BLUE     = RGBColor(0x27, 0x6D, 0xC8)

SS  = r"C:\Users\rkvlu\Desktop\MyTunes\screenshots"
OUT = r"C:\Users\rkvlu\Desktop\MyTunes\MyTunes_Presentation_Final.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, colour=BG_DARK):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = colour

def rect(slide, l, t, w, h, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(lw or 1)
    else:    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=13, bold=False, colour=WHITE,
        align=PP_ALIGN.LEFT, font="Segoe UI"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = colour; r.font.name = font
    return tb

def header(slide, title, subtitle="", accent=ACCENT):
    rect(slide, 0, 0, 13.33, 1.28, fill=BG_CARD)
    rect(slide, 0, 1.24, 13.33, 0.04, fill=accent)
    txt(slide, title, 0.4, 0.1, 12, 0.62, size=28, bold=True, colour=accent)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.72, 12, 0.38, size=12, colour=GREY)

def pill(slide, l, t, w, h, label, bg_col, tc=WHITE, sz=10):
    rect(slide, l, t, w, h, fill=bg_col)
    txt(slide, label, l, t + (h-0.18)/2, w, 0.22,
        size=sz, bold=True, colour=tc, align=PP_ALIGN.CENTER)

def bullet_block(slide, icon_col, title_text, body_text, l, t, w):
    """One bullet: coloured bar | bold title | grey body."""
    rect(slide, l, t + 0.04, 0.06, 0.26, fill=icon_col)
    txt(slide, title_text, l + 0.2, t, w - 0.25, 0.32,
        size=12, bold=True, colour=WHITE)
    txt(slide, body_text,  l + 0.2, t + 0.32, w - 0.25, 0.42,
        size=11, colour=GREY)

def picture(slide, path, l, t, w, h):
    from PIL import Image as PILImage
    img = PILImage.open(path)
    iw, ih = img.size
    if (iw/ih) > (w/h): fit_w=w; fit_h=w/(iw/ih)
    else:                fit_h=h; fit_w=h*(iw/ih)
    ol = l + (w-fit_w)/2; ot = t + (h-fit_h)/2
    slide.shapes.add_picture(path, Inches(ol), Inches(ot),
                             Inches(fit_w), Inches(fit_h))

def screenshot_slide(title, subtitle, img_path, points, accent_col=ACCENT2):
    s = prs.slides.add_slide(blank); bg(s)
    header(s, title, subtitle, accent=accent_col)
    SL,ST,SW,SH = 0.35, 1.38, 6.85, 5.98
    rect(s, SL-0.04, ST-0.04, SW+0.08, SH+0.08,
         fill=RGBColor(0x1A,0x1A,0x2A), line=accent_col, lw=1.5)
    rect(s, SL, ST, SW, 0.3, fill=RGBColor(0x2D,0x2D,0x2D))
    for i, col in enumerate([RED, YELLOW, ACCENT]):
        c = s.shapes.add_shape(9, Inches(SL+0.1+i*0.18), Inches(ST+0.08),
                               Inches(0.12), Inches(0.12))
        c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()
    txt(s, subtitle.split("|")[0].strip() if "|" in subtitle else subtitle,
        SL+0.6, ST+0.05, SW-0.65, 0.22, size=8, colour=GREY)
    if img_path and os.path.exists(img_path):
        picture(s, img_path, SL, ST+0.3, SW, SH-0.3)
    else:
        rect(s, SL, ST+0.3, SW, SH-0.3, fill=RGBColor(0x18,0x18,0x18))
        txt(s, "[Screenshot]", SL+2.5, ST+3, 2, 0.4,
            size=12, colour=GREY, align=PP_ALIGN.CENTER)
    EL,ET,EW,EH = 7.3, 1.38, 5.85, 5.98
    rect(s, EL, ET, EW, EH, fill=BG_CARD,
         line=RGBColor(0x2A,0x2A,0x3A), lw=0.5)
    txt(s, "How it Works", EL+0.25, ET+0.18, EW-0.4, 0.4,
        size=15, bold=True, colour=accent_col)
    rect(s, EL+0.25, ET+0.6, EW-0.5, 0.03, fill=accent_col)
    y = ET + 0.72
    for tp, bp in points:
        rect(s, EL+0.22, y+0.04, 0.06, 0.26, fill=accent_col)
        txt(s, tp, EL+0.4, y, EW-0.6, 0.32, size=11, bold=True, colour=WHITE)
        txt(s, bp, EL+0.4, y+0.33, EW-0.6, 0.42, size=10, colour=GREY)
        y += 0.82


# =============================================================================
# SLIDE 1 – Title
# =============================================================================
s1 = prs.slides.add_slide(blank); bg(s1)
rect(s1, 0, 4.85, 13.33, 0.07, fill=ACCENT)
rect(s1, 0, 4.92, 13.33, 0.04, fill=ACCENT2)
txt(s1, "🎵", 9.6, 0.5, 3.2, 3.2, size=110, align=PP_ALIGN.CENTER)
txt(s1, "MyTunes",  0.6, 0.9, 8.8, 1.2, size=66, bold=True, colour=ACCENT)
txt(s1, "Ad-Free Social Music Streaming Platform",
    0.6, 2.1, 9, 0.55, size=20, colour=GREY)
rect(s1, 0.6, 2.82, 5.2, 0.05, fill=ACCENT2)
rect(s1, 0.5, 3.05, 7.8, 3.85, fill=BG_CARD,
     line=RGBColor(0x33,0x33,0x44), lw=0.5)
txt(s1, "PROJECT TEAM", 0.8, 3.22, 7, 0.32, size=10, bold=True, colour=ACCENT2)
for i, (role, name) in enumerate([
    ("Team Member 1","[ Full Name ]"),
    ("Team Member 2","[ Full Name ]"),
    ("Team Member 3","[ Full Name ]"),
    ("Team Member 4","[ Full Name ]"),
]):
    y = 3.65 + i*0.7
    pill(s1, 0.8, y, 1.7, 0.34, role, ACCENT2, sz=9)
    txt(s1, name, 2.6, y+0.01, 5.3, 0.34, size=14, bold=True, colour=WHITE)
txt(s1, "Department of Computer Science  |  Academic Year 2024-25",
    0.6, 7.15, 12, 0.28, size=10, colour=GREY)


# =============================================================================
# SLIDE 2 – Introduction
# =============================================================================
s2 = prs.slides.add_slide(blank); bg(s2)
header(s2, "Introduction", "What is MyTunes?", accent=ACCENT)

# Left column – what it is
rect(s2, 0.35, 1.42, 6.2, 5.82, fill=BG_CARD,
     line=ACCENT, lw=0.8)
txt(s2, "About the Project", 0.6, 1.55, 5.8, 0.38,
    size=14, bold=True, colour=ACCENT)
rect(s2, 0.55, 1.96, 5.7, 0.03, fill=ACCENT)

intro_points = [
    ("What is MyTunes?",
     "A free, ad-free music streaming web app that lets users listen to music, "
     "build playlists, and connect with friends – all in one place."),
    ("Why we built it?",
     "Popular platforms like Spotify and Apple Music charge for ad-free listening "
     "and lock social features behind premium plans. MyTunes gives everyone these "
     "features for free."),
    ("Who is it for?",
     "Music lovers who want a free, social streaming experience without ads "
     "or subscription fees."),
    ("How it works?",
     "Users sign up, search for songs, create playlists, like tracks, add friends, "
     "and can even listen to the same song together in real time."),
]
y = 2.08
for btitle, bbody in intro_points:
    rect(s2, 0.55, y+0.04, 0.06, 0.25, fill=ACCENT)
    txt(s2, btitle, 0.75, y,    5.5, 0.3,  size=11, bold=True, colour=WHITE)
    txt(s2, bbody,  0.75, y+0.3, 5.5, 0.55, size=10, colour=GREY)
    y += 0.95

# Right column – key highlights
rect(s2, 6.75, 1.42, 6.25, 5.82, fill=BG_CARD,
     line=ACCENT2, lw=0.8)
txt(s2, "Key Highlights", 7.0, 1.55, 5.8, 0.38,
    size=14, bold=True, colour=ACCENT2)
rect(s2, 6.95, 1.96, 5.7, 0.03, fill=ACCENT2)

highlights = [
    (ACCENT,  "Free & Ad-Free",       "No ads, no subscriptions – ever"),
    (ACCENT2, "Social Platform",      "Add friends, share music, see what they listen to"),
    (TEAL,    "Listen Together",      "Real-time sync – same song, same moment, different places"),
    (YELLOW,  "Mood Discovery",       "Find music based on your current mood"),
    (RED,     "Secure & Private",     "JWT auth, bcrypt passwords, private playlist options"),
    (BLUE,    "Full-Stack Project",   "Node.js + MongoDB backend + Vanilla JS frontend"),
]
y = 2.08
for col, htitle, hbody in highlights:
    rect(s2, 6.95, y, 5.7, 0.75, fill=RGBColor(0x16,0x16,0x26),
         line=col, lw=0.6)
    rect(s2, 6.95, y, 0.06, 0.75, fill=col)
    txt(s2, htitle, 7.12, y+0.08, 5.3, 0.28, size=11, bold=True, colour=col)
    txt(s2, hbody,  7.12, y+0.38, 5.3, 0.28, size=10, colour=GREY)
    y += 0.87


# =============================================================================
# SLIDE 3 – Problem Statement
# =============================================================================
s3 = prs.slides.add_slide(blank); bg(s3)
header(s3, "Problem Statement", "The gap MyTunes fills", accent=RED)

# Central problem statement box
rect(s3, 0.5, 1.42, 12.35, 1.0, fill=RGBColor(0x2A,0x0A,0x0A),
     line=RED, lw=1.5)
txt(s3,
    "Music lovers today are forced to either tolerate constant ads on free tiers OR pay expensive monthly "
    "subscriptions just to enjoy ad-free streaming. Social and collaborative features remain locked behind "
    "premium plans – making music a solo, costly experience.",
    0.75, 1.52, 11.85, 0.82, size=12, colour=WHITE)

problems = [
    (RED,    "Costly Ad-Free Streaming",
             "Spotify Free = ads every 2-3 songs. Ad-free requires ₹119/month. "
             "Many students and users cannot afford premium plans."),
    (YELLOW, "No Free Social Features",
             "Following friends, seeing what they listen to, and sharing music requires "
             "paid tiers on most platforms. Free users are isolated."),
    (ACCENT2,"No Real-Time Group Listening (Free)",
             "Listening to the same song together with friends in real time is not available "
             "on any major free streaming platform without third-party tools."),
    (TEAL,   "Limited Mood-Based Discovery",
             "Free tiers offer very limited playlist and mood-based browsing. "
             "Personalised recommendations are gated behind subscriptions."),
    (ACCENT, "Complex Privacy Controls",
             "Users cannot easily control who sees their playlists. "
             "Most platforms lack simple Public / Friends-only / Private options on free plans."),
    (BLUE,   "No Unified Platform",
             "Users juggle multiple apps for streaming, social features, and group listening. "
             "There is no single free platform that combines all of these."),
]

cols = 3
cw, ch = 4.1, 1.55
for idx, (col, ptitle, pbody) in enumerate(problems):
    row, c = divmod(idx, cols)
    l = 0.4 + c * (cw + 0.12)
    t = 2.58 + row * (ch + 0.15)
    rect(s3, l, t, cw, ch, fill=BG_CARD, line=col, lw=1.0)
    rect(s3, l, t, 0.07, ch, fill=col)
    txt(s3, ptitle, l+0.18, t+0.1, cw-0.25, 0.35, size=11, bold=True, colour=col)
    txt(s3, pbody,  l+0.18, t+0.5, cw-0.25, 0.92, size=10, colour=GREY)


# =============================================================================
# SLIDE 4 – Objectives
# =============================================================================
s4 = prs.slides.add_slide(blank); bg(s4)
header(s4, "Objectives", "What MyTunes aims to achieve", accent=YELLOW)

# Left: main objectives
rect(s4, 0.35, 1.42, 6.2, 5.82, fill=BG_CARD, line=YELLOW, lw=0.8)
txt(s4, "Primary Objectives", 0.58, 1.55, 5.8, 0.38,
    size=14, bold=True, colour=YELLOW)
rect(s4, 0.55, 1.96, 5.7, 0.03, fill=YELLOW)

obj_main = [
    ("Build a Free, Ad-Free Platform",
     "Create a music streaming app that requires no subscription and shows zero ads to users."),
    ("Implement Secure Authentication",
     "Use JWT tokens and bcrypt hashing so user accounts and passwords are always protected."),
    ("Enable Playlist Management",
     "Let users create, edit, and delete playlists with Public, Friends-only, or Private visibility."),
    ("Real-Time Listen Together",
     "Build a Socket.io-powered feature so friends can listen to the same song simultaneously."),
    ("Social / Friends System",
     "Allow users to send friend requests, see friends' activity, and discover music socially."),
]
y = 2.08
for otitle, obody in obj_main:
    rect(s4, 0.55, y+0.04, 0.06, 0.25, fill=YELLOW)
    txt(s4, otitle, 0.75, y,     5.5, 0.3,  size=11, bold=True, colour=WHITE)
    txt(s4, obody,  0.75, y+0.3, 5.5, 0.52, size=10, colour=GREY)
    y += 0.9

# Right: secondary objectives
rect(s4, 6.75, 1.42, 6.25, 5.82, fill=BG_CARD, line=ACCENT, lw=0.8)
txt(s4, "Secondary Objectives", 7.0, 1.55, 5.8, 0.38,
    size=14, bold=True, colour=ACCENT)
rect(s4, 6.95, 1.96, 5.7, 0.03, fill=ACCENT)

obj_sec = [
    ("Mood-Based Music Discovery",
     "Tag songs by mood (Happy, Sad, Energetic, Calm) and let users browse by feeling."),
    ("User Profiles & Stats",
     "Show each user's top tracks, total plays, public playlists, and friends count."),
    ("Automated Testing",
     "Write Jest + Supertest test suites to verify all API routes work correctly."),
    ("Clean & Responsive UI",
     "Build a dark-themed, Spotify-inspired interface that works well at 1440px wide."),
    ("Structured REST API",
     "Design clear, consistent API endpoints with proper HTTP status codes and error messages."),
    ("Logging & Monitoring",
     "Use Winston logger to track server events, errors, and API requests for debugging."),
]
y = 2.08
for otitle, obody in obj_sec:
    rect(s4, 6.95, y+0.04, 0.06, 0.25, fill=ACCENT)
    txt(s4, otitle, 7.15, y,     5.5, 0.3,  size=11, bold=True, colour=WHITE)
    txt(s4, obody,  7.15, y+0.3, 5.5, 0.52, size=10, colour=GREY)
    y += 0.9


# =============================================================================
# SLIDE 5 – Scope of the Project
# =============================================================================
s5 = prs.slides.add_slide(blank); bg(s5)
header(s5, "Scope of the Project", "What is included and what is not", accent=TEAL)

# In Scope (left)
rect(s5, 0.35, 1.42, 6.0, 5.82, fill=BG_CARD,
     line=ACCENT, lw=0.8)
txt(s5, "✅  IN SCOPE", 0.6, 1.52, 5.6, 0.38,
    size=14, bold=True, colour=ACCENT)
rect(s5, 0.55, 1.93, 5.6, 0.03, fill=ACCENT)

in_scope = [
    "User Registration & Login (Email + Password with JWT)",
    "Music Search – search songs by title or artist name",
    "Full audio player with play, pause, skip, shuffle, repeat, volume",
    "Playlist creation, editing, and deletion with visibility controls",
    "Liked Songs – save and manage favourite tracks",
    "Mood-Based Browsing – filter songs by mood tags",
    "Friends System – send, accept, decline friend requests",
    "Listen Together – real-time group sessions via Socket.io",
    "User Profile – display name, bio, stats, public playlists",
    "Queue management – add songs to a persistent play queue",
    "Toast notifications and context menus for song actions",
    "Automated backend tests with Jest + Supertest",
]
y = 2.05
for item in in_scope:
    rect(s5, 0.55, y+0.06, 0.08, 0.08, fill=ACCENT)
    txt(s5, item, 0.76, y, 5.35, 0.32, size=10, colour=GREY)
    y += 0.36

# Out of Scope (right)
rect(s5, 6.65, 1.42, 6.3, 5.82, fill=BG_CARD,
     line=RED, lw=0.8)
txt(s5, "❌  OUT OF SCOPE", 6.9, 1.52, 5.8, 0.38,
    size=14, bold=True, colour=RED)
rect(s5, 6.85, 1.93, 5.8, 0.03, fill=RED)

out_scope = [
    ("Mobile Application", "Web app only – no iOS/Android app this semester"),
    ("Offline Mode / Downloads", "All streaming requires an active internet connection"),
    ("ML Recommendations", "No AI-based 'you might also like' engine yet"),
    ("Artist / Creator Accounts", "Only listener accounts; artists cannot upload directly"),
    ("Payment / Subscription", "No premium tier or payment gateway integration"),
    ("In-App Chat", "Friends can see activity but cannot send messages yet"),
    ("Live Radio / Broadcasting", "Not a live streaming or radio service"),
    ("Desktop / Mobile App", "Not packaged as an Electron or PWA offline app"),
]
y = 2.05
for otitle, obody in out_scope:
    rect(s5, 6.85, y+0.04, 0.06, 0.24, fill=RED)
    txt(s5, otitle, 7.06, y,     5.65, 0.28, size=11, bold=True, colour=WHITE)
    txt(s5, obody,  7.06, y+0.3, 5.65, 0.34, size=10, colour=GREY)
    y += 0.72


# =============================================================================
# SLIDE 6 – Technologies Used
# =============================================================================
s6 = prs.slides.add_slide(blank); bg(s6)
header(s6, "Technologies Used", "The tech stack that powers MyTunes", accent=BLUE)

tech_sections = [
    ("Frontend", ACCENT2, [
        ("HTML5 + CSS3",      "Page structure and dark-theme styling with CSS variables and Grid/Flexbox"),
        ("JavaScript (ES6+)", "Vanilla JS SPA – no framework. Handles routing, API calls, DOM updates"),
        ("Font Awesome 6",    "Icon library for all UI icons (play, heart, users, etc.)"),
        ("Socket.io Client",  "Connects to the server's WebSocket for real-time sync features"),
    ]),
    ("Backend", ACCENT, [
        ("Node.js",           "JavaScript runtime for the server – handles all requests asynchronously"),
        ("Express.js",        "Web framework for defining REST API routes and middleware"),
        ("Socket.io",         "Real-time bidirectional communication for Listen Together sessions"),
        ("Multer",            "Handles file uploads (profile pictures, cover images)"),
    ]),
    ("Database & Security", TEAL, [
        ("MongoDB",           "NoSQL document database – stores Users, Songs, Playlists, Sessions"),
        ("Mongoose ODM",      "Schema modeling for MongoDB with validation and relationships"),
        ("JWT",               "JSON Web Tokens for stateless authentication (24h access + 30d refresh)"),
        ("bcrypt.js",         "Hashes passwords before storing – even the admin cannot see passwords"),
    ]),
    ("Tools & Testing", YELLOW, [
        ("Jest + Supertest",  "Automated test suite for all API endpoints and controller logic"),
        ("Winston",           "Structured logging of server events, errors, and API requests"),
        ("Helmet",            "Sets security HTTP headers to protect against common web attacks"),
        ("Deezer API",        "External music API used to fetch song metadata, covers, and previews"),
    ]),
]

cw, ch_sec = 3.12, 3.3
for col_idx, (sec_title, col, items) in enumerate(tech_sections):
    lx = 0.3 + col_idx * (cw + 0.1)
    ty = 1.42
    rect(s6, lx, ty, cw, ch_sec + 0.05, fill=BG_CARD, line=col, lw=1.2)
    rect(s6, lx, ty, cw, 0.42, fill=col)
    txt(s6, sec_title, lx, ty+0.07, cw, 0.3,
        size=12, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    iy = ty + 0.52
    for tech, desc in items:
        rect(s6, lx+0.12, iy+0.04, 0.06, 0.2, fill=col)
        txt(s6, tech, lx+0.28, iy,      cw-0.38, 0.26, size=10, bold=True, colour=WHITE)
        txt(s6, desc, lx+0.28, iy+0.26, cw-0.38, 0.35, size=9, colour=GREY)
        iy += 0.65

# Why these choices – bottom strip
rect(s6, 0.3, 5.0, 12.73, 2.28, fill=BG_CARD,
     line=RGBColor(0x2A,0x2A,0x3A), lw=0.5)
txt(s6, "Why This Stack?", 0.55, 5.08, 12, 0.35, size=13, bold=True, colour=BLUE)
reasons = [
    ("Node.js + Express", "Fast, lightweight, perfect for API servers. JavaScript everywhere (front + back)."),
    ("MongoDB",           "Schema-flexible – easy to add new fields without migrations. Great for rapid development."),
    ("Socket.io",         "Simplest way to add real-time features; handles WebSocket fallbacks automatically."),
    ("Vanilla JS",        "No build step, no framework overhead – great for learning fundamentals and fast deployment."),
]
rx = 0.5
for rtitle, rbody in reasons:
    rect(s6, rx, 5.5, 0.06, 0.28, fill=BLUE)
    txt(s6, rtitle, rx+0.14, 5.46, 2.8, 0.28, size=10, bold=True, colour=WHITE)
    txt(s6, rbody,  rx+0.14, 5.74, 2.8, 0.42, size=9, colour=GREY)
    rx += 3.15


# =============================================================================
# SLIDE 7 – Project Overview (existing content)
# =============================================================================
s7 = prs.slides.add_slide(blank); bg(s7)
header(s7, "Project Overview", "Feature summary", accent=ACCENT2)

txt(s7,
    "MyTunes is a full-stack, ad-free social music streaming web application built with "
    "Node.js/Express (backend) and vanilla HTML/CSS/JS (frontend). It delivers a Spotify-like "
    "experience while adding social features – friends, playlist sharing, mood-based discovery, "
    "and real-time synchronized group listening via Socket.io.",
    0.4, 1.38, 12.5, 0.95, size=13, colour=GREY)

feature_cards = [
    (ACCENT,  "🔐  Auth & Security",   "JWT tokens · bcrypt hashing\nRemember-Me refresh tokens"),
    (ACCENT2, "🎵  Music Streaming",   "Browse, search & stream songs\nFull-featured audio player"),
    (TEAL,    "📋  Playlists",          "Create, edit & manage playlists\nPublic / Friends / Private"),
    (YELLOW,  "😊  Mood Discovery",    "Browse tracks by mood category\nHappy, Sad, Energetic…"),
    (RED,     "👥  Friends & Social",  "Send requests · Friend activity\nShare music in real time"),
    (RGBColor(0x8E,0x44,0xAD),"🎧  Listen Together",
                                        "Socket.io room sync\nGroup listening sessions"),
]
cw2, ch2 = 4.0, 1.55
for idx,(col,ctitle,desc) in enumerate(feature_cards):
    row, c = divmod(idx, 3)
    l = 0.3 + c*(cw2+0.22); t = 2.45 + row*(ch2+0.18)
    rect(s7, l, t, cw2, ch2, fill=BG_CARD, line=col, lw=1.2)
    rect(s7, l, t, 0.07, ch2, fill=col)
    txt(s7, ctitle, l+0.18, t+0.1,  cw2-0.25, 0.38, size=12, bold=True, colour=col)
    txt(s7, desc,   l+0.18, t+0.52, cw2-0.25, 0.85, size=11, colour=GREY)

txt(s7, "Tech Stack:", 0.4, 7.12, 1.4, 0.28, size=10, bold=True, colour=WHITE)
x = 1.75
for label, col in [("Node.js",ACCENT),("Express",ACCENT2),("MongoDB",TEAL),
                   ("Socket.io",YELLOW),("JWT",RED),("HTML / CSS / JS",BLUE)]:
    pill(s7, x, 7.12, max(1.3, len(label)*0.12), 0.27, label, col, sz=9)
    x += max(1.3, len(label)*0.12) + 0.1


# =============================================================================
# SLIDES 8–12 – Screenshot slides (Login, Register, Dashboard, etc.)
# =============================================================================
screenshot_slide(
    "Page 1 – Login & Register",
    "localhost:5500/index.html  |  Entry point",
    os.path.join(SS, "01_login.png"),
    [("Split-Layout Design","Left hero panel shows branding; right card holds Login/Register forms."),
     ("Login Flow","Email + Password → POST /api/auth/login → JWT stored → redirect to dashboard."),
     ("Register Flow","Username, Email, Password → POST /api/auth/register → auto-login on success."),
     ("Remember Me","Issues a 30-day refresh token so users stay signed in across sessions."),
     ("Password Toggle","Eye icon switches the password field between hidden and visible text."),
     ("Route Guard","If a valid JWT already exists the page skips login and opens the dashboard."),
    ], accent_col=ACCENT)

screenshot_slide(
    "Page 1b – Register Form",
    "localhost:5500/index.html  |  Register tab",
    os.path.join(SS, "01b_register.png"),
    [("Username Field","Must be unique. Validated on server via express-validator before saving."),
     ("Email Validation","Format checked on client and server. Duplicate email returns 409 Conflict."),
     ("Password Rules","Minimum strength rules enforced by validator; bcrypt hashes before storing."),
     ("Confirm Password","Client check ensures both fields match before the form is submitted."),
     ("Auto-Login","Successful registration returns a JWT and takes user straight to dashboard."),
     ("Toggle Back","Login / Register buttons switch between the two forms with a CSS transition."),
    ], accent_col=ACCENT)

screenshot_slide(
    "Page 2 – Dashboard: Home",
    "localhost:5500/dashboard.html  |  Main hub",
    os.path.join(SS, "02_dashboard_home.png"),
    [("3-Column Layout","Sidebar (nav) | Scrollable content area | Optional Queue/Friends panel."),
     ("Time-Aware Greeting","Shows Good Morning/Afternoon/Evening based on the user's local time."),
     ("Trending Songs","Fetched from GET /api/music/trending. Shown as clickable song cards."),
     ("Browse by Mood","6 mood tiles. Each loads GET /api/music/mood?tag=<mood> when clicked."),
     ("Your Playlists","Loaded via GET /api/playlists/mine. Click a card to open the full playlist."),
     ("Persistent Player","Fixed footer bar: art, title, play/pause, next, prev, shuffle, volume."),
    ], accent_col=ACCENT2)

screenshot_slide(
    "Page 3 – Search",
    "dashboard.html  |  Search section",
    os.path.join(SS, "03_search.png"),
    [("Real-time Search","Debounced input hits GET /api/music/search?q=<query> as the user types."),
     ("Song Cards","Each result shows thumbnail, title, artist, duration. Click to play instantly."),
     ("Like from Results","Heart icon toggles POST /api/music/:id/like. State syncs to Liked Songs."),
     ("Add to Playlist","'+' button opens the Add to Playlist modal. Song added via PATCH /api/playlists/:id."),
     ("Context Menu","Right-click → Play, Like, Add to Playlist, Add to Queue, Share."),
     ("Queue Integration","Songs added to queue appear in the right panel and play in order."),
    ], accent_col=TEAL)

screenshot_slide(
    "Page 4 – Library",
    "dashboard.html  |  Library section",
    os.path.join(SS, "04_library.png"),
    [("Playlists Grid","All user playlists as cards. Create new ones with the + Create button."),
     ("Recently Played","Tracks stored in PlayHistory. Fetched from GET /api/music/history."),
     ("Create Playlist","Modal: Name, Description, Visibility (Public/Friends/Private)."),
     ("Playlist Card Click","Opens full detail view with song list, hero image, play/shuffle."),
     ("Sidebar Playlists","Playlists also listed in the sidebar for quick access from any page."),
     ("Delete Playlist","Trash icon → DELETE /api/playlists/:id after a confirm dialog."),
    ], accent_col=ACCENT2)

screenshot_slide(
    "Page 5 – Liked Songs",
    "dashboard.html  |  Liked Songs",
    os.path.join(SS, "05_liked_songs.png"),
    [("Saved Songs","All liked songs stored here. Fetched via GET /api/music/liked."),
     ("Play All","Green play button loads all liked songs into the queue from song #1."),
     ("Shuffle","Randomises the playback order of the liked songs queue."),
     ("Song Table","Columns: #, Title, Artist, Duration, Like toggle. Click row to play."),
     ("Unlike","Click the heart again to unlike. POST /api/music/:id/like toggles state."),
     ("Song Count","Total count shown in header, updated dynamically after each like/unlike."),
    ], accent_col=RED)

screenshot_slide(
    "Page 6 – Friends",
    "dashboard.html  |  Friends section",
    os.path.join(SS, "06_friends.png"),
    [("Search Users","Type a username or email to find users → POST /api/friends/search."),
     ("Send Request","Click Add → POST /api/friends/request. Recipient sees it in Pending."),
     ("Accept / Decline","POST /api/friends/accept or /decline. Creates a mutual Friendship."),
     ("Friends List","Shows confirmed friends. Active friends show their current song."),
     ("Real-time Activity","Right panel updates live via Socket.io 'user-listening' events."),
     ("Remove Friend","Three-dot menu → Remove Friend. Both users' lists update instantly."),
    ], accent_col=TEAL)

screenshot_slide(
    "Page 7 – Listen Together",
    "dashboard.html  |  Listen Together",
    os.path.join(SS, "07_listen_together.png"),
    [("Create a Session","Host names it, sets visibility. Server generates a unique 6-char code."),
     ("Join by Code","Others enter the code in Join modal → POST /api/sessions/join."),
     ("Socket.io Sync","Play, Pause, Seek, Skip events broadcast to all room members instantly."),
     ("Host Control","Only the host controls the song. All listeners follow host's state."),
     ("Live Member List","Shows who is in the session, updated as people join or leave."),
     ("Leave Session","Any member can leave. If host leaves, session ends for everyone."),
    ], accent_col=BLUE)

screenshot_slide(
    "Page 8 – Profile",
    "dashboard.html  |  Profile section",
    os.path.join(SS, "08_profile.png"),
    [("Profile Hero","Shows avatar, display name, total plays, playlists, and friends count."),
     ("Edit Profile","Modal to update Display Name, Bio, picture. PATCH /api/users/profile."),
     ("Top Tracks","Most-played songs from PlayHistory, sorted by play count descending."),
     ("Public Playlists","Playlists marked Public are visible to other users."),
     ("Stats Panel","Plays, playlists, friends counts from GET /api/users/me + history."),
     ("Settings","Dark Mode, AutoPlay, Language preferences stored in user document."),
    ], accent_col=ACCENT)

screenshot_slide(
    "Feature – Create Playlist Modal",
    "Modal overlay from Library or Home page",
    os.path.join(SS, "09_create_playlist.png"),
    [("Playlist Name","Required. Validated client-side and server-side before saving."),
     ("Description","Optional text describing the playlist. Stored in MongoDB document."),
     ("Visibility Options","Public – anyone. Friends Only – only friends. Private – only you."),
     ("Create Button","POST /api/playlists → New playlist added to sidebar and Library instantly."),
     ("Cancel / Escape","Dismisses modal without creating. Escape key also closes it."),
     ("Instant Refresh","Playlists grid and sidebar update without a full page reload."),
    ], accent_col=ACCENT2)


# =============================================================================
# SLIDE – System Architecture
# =============================================================================
sa = prs.slides.add_slide(blank); bg(sa)
header(sa, "System Architecture", "How all layers connect", accent=BLUE)

layers = [
    ("FRONTEND",  "HTML · CSS · JS\nVanilla SPA\n\n• index.html\n• dashboard.html\n• playlist.html",
     ACCENT2, 0.35),
    ("REST API",  "Node.js + Express\n\n• /api/auth\n• /api/music\n• /api/playlists\n• /api/friends\n• /api/sessions",
     ACCENT, 3.75),
    ("REAL-TIME", "Socket.io\n\n• Room events:\n  play/pause/seek\n• Friends activity\n  broadcast",
     BLUE, 7.15),
    ("DATABASE",  "MongoDB\n+ Mongoose ODM\n\n• User · Song\n• Playlist\n• Friendship\n• ListeningSession",
     YELLOW, 10.55),
]
for title, body, col, lx in layers:
    rect(sa, lx, 1.42, 2.75, 4.95, fill=BG_CARD, line=col, lw=1.5)
    rect(sa, lx, 1.42, 2.75, 0.5, fill=col)
    txt(sa, title, lx, 1.49, 2.75, 0.36,
        size=12, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    txt(sa, body, lx+0.2, 2.06, 2.38, 3.95, size=10, colour=GREY)

for ax in [3.1, 6.5, 9.9]:
    txt(sa, "→", ax+0.06, 3.6, 0.6, 0.4, size=18, colour=GREY, align=PP_ALIGN.CENTER)

badges = [
    ("🔐 Security",   "JWT 24h + 30d Refresh · bcrypt hashing · Protected route middleware", ACCENT),
    ("🧪 Testing",    "Jest + Supertest · Unit & Integration tests for controllers & routes",  ACCENT2),
    ("🚀 Deployment", "Backend: Node.js port 5000 · Frontend: static files via Express/LiveServer", TEAL),
]
for i,(btitle,bbody,bcol) in enumerate(badges):
    bl = 0.35 + i*4.35
    rect(sa, bl, 6.55, 4.1, 0.72, fill=BG_CARD, line=bcol, lw=0.8)
    txt(sa, btitle, bl+0.15, 6.6,  3.85, 0.28, size=11, bold=True, colour=bcol)
    txt(sa, bbody,  bl+0.15, 6.88, 3.85, 0.28, size=9, colour=GREY)


# =============================================================================
# SLIDE – Outcomes / Results
# =============================================================================
or_ = prs.slides.add_slide(blank); bg(or_)
header(or_, "Outcomes & Results", "What we achieved by the end of this semester", accent=ACCENT)

# Top summary bar
rect(or_, 0.35, 1.42, 12.63, 0.82, fill=RGBColor(0x0A,0x2A,0x0A),
     line=ACCENT, lw=1.2)
txt(or_,
    "MyTunes was successfully built and deployed as a fully working full-stack web application. "
    "All planned features for this semester were implemented, tested, and are functional.",
    0.6, 1.52, 12.1, 0.65, size=12, colour=WHITE)

# Stat boxes
stats = [
    ("8+",  "Major Features\nDelivered", ACCENT),
    ("20+", "REST API\nEndpoints", ACCENT2),
    ("100%","Core Features\nWorking",    TEAL),
    ("Jest","Automated\nTests Written",  YELLOW),
]
for i,(num, label, col) in enumerate(stats):
    bx = 0.35 + i * 3.2
    rect(or_, bx, 2.38, 3.0, 1.2, fill=BG_CARD, line=col, lw=1.2)
    txt(or_, num,   bx, 2.45, 3.0, 0.65,
        size=32, bold=True, colour=col, align=PP_ALIGN.CENTER)
    txt(or_, label, bx, 3.08, 3.0, 0.45,
        size=10, colour=GREY, align=PP_ALIGN.CENTER)

# Outcome bullets (2 columns)
outcomes_left = [
    ("Secure Authentication Working",
     "Users can register, login, and stay logged in with JWT + refresh tokens. Passwords are bcrypt-hashed."),
    ("Music Streaming Functional",
     "Songs load from the Deezer API. The audio player supports play, pause, skip, shuffle, repeat, and volume."),
    ("Playlists Fully Operational",
     "Users can create, rename, delete playlists and add/remove songs. Visibility controls work correctly."),
    ("Mood Discovery Works",
     "6 mood categories return correct filtered song lists from the backend based on mood tags."),
]
outcomes_right = [
    ("Friends System Complete",
     "Friend requests, acceptance, declination, and removal all work. Friends' listening activity is visible."),
    ("Listen Together Functional",
     "Multiple users can join a session and listen to the same song in sync via Socket.io rooms."),
    ("Profile Page Live",
     "User profile shows real stats: play count, playlists, friends. Editing display name and bio works."),
    ("Test Suite Passes",
     "Jest + Supertest tests cover auth routes, music routes, and playlist operations. All passing."),
]
y = 3.72
for (lt, lb), (rt, rb) in zip(outcomes_left, outcomes_right):
    rect(or_, 0.35, y+0.04, 0.06, 0.25, fill=ACCENT)
    txt(or_, lt, 0.55, y,     6.0, 0.28, size=11, bold=True, colour=WHITE)
    txt(or_, lb, 0.55, y+0.3, 6.0, 0.38, size=10, colour=GREY)
    rect(or_, 6.7, y+0.04, 0.06, 0.25, fill=ACCENT2)
    txt(or_, rt, 6.9, y,     6.0, 0.28, size=11, bold=True, colour=WHITE)
    txt(or_, rb, 6.9, y+0.3, 6.0, 0.38, size=10, colour=GREY)
    y += 0.82


# =============================================================================
# SLIDE – Future Work / Next Semester Plan
# =============================================================================
fw = prs.slides.add_slide(blank); bg(fw)
header(fw, "Future Work", "Next Semester Implementation Plan", accent=ACCENT2)

txt(fw,
    "MyTunes has a strong foundation. In the next semester, we plan to expand it with advanced features "
    "that make it even more powerful, social, and accessible to users on all devices.",
    0.4, 1.38, 12.5, 0.72, size=12, colour=GREY)

future_items = [
    (ACCENT,  "📱 Mobile Application",
              "Phase 1 Priority",
              "Build a React Native mobile app for iOS and Android. "
              "All existing features (streaming, playlists, friends, Listen Together) available on phone."),
    (ACCENT2, "🤖 ML Recommendations",
              "Phase 2",
              "Train a recommendation model on user play history to suggest 'You might also like' songs. "
              "Personalise the Home page feed for each user."),
    (BLUE,    "💬 In-App Chat",
              "Phase 1 Priority",
              "Add a real-time chat system so friends can message each other. "
              "Send song links directly in chat with a one-click play button."),
    (TEAL,    "⬇️ Offline Mode",
              "Phase 2",
              "Allow users to download songs for offline playback using IndexedDB or a PWA service worker. "
              "Syncs back when the connection is restored."),
    (YELLOW,  "🎤 Artist Accounts",
              "Phase 2",
              "Let artists/creators register a special account, upload their own music, "
              "and view streams and listener statistics for their tracks."),
    (RED,     "📊 Analytics Dashboard",
              "Phase 1 Priority",
              "Give users a personal stats page: listening time per day, top genres, "
              "most-played hours, and monthly listening trends with charts."),
]

cw3, ch3 = 4.08, 1.68
for idx,(col, ftitle, phase, fbody) in enumerate(future_items):
    row, c = divmod(idx, 3)
    l = 0.3 + c*(cw3+0.08)
    t = 2.2 + row*(ch3+0.16)
    rect(fw, l, t, cw3, ch3, fill=BG_CARD, line=col, lw=1.0)
    rect(fw, l, t, cw3, 0.38, fill=col)
    txt(fw, ftitle, l+0.12, t+0.06, cw3-0.2, 0.28, size=11, bold=True, colour=WHITE)
    # Phase badge
    rect(fw, l+cw3-1.15, t+0.42, 1.1, 0.25, fill=RGBColor(0x10,0x10,0x20))
    txt(fw, phase, l+cw3-1.15, t+0.44, 1.1, 0.22,
        size=8, colour=col, align=PP_ALIGN.CENTER)
    txt(fw, fbody, l+0.12, t+0.72, cw3-0.22, 0.85, size=10, colour=GREY)

# Bottom note
rect(fw, 0.3, 7.05, 12.73, 0.36, fill=BG_CARD)
txt(fw,
    "Phase 1 features will be completed by mid-next semester.  "
    "Phase 2 features will be research + prototyped by end of next semester.",
    0.55, 7.1, 12.2, 0.28, size=10, colour=GREY)


# =============================================================================
# SLIDE – Thank You
# =============================================================================
ty_ = prs.slides.add_slide(blank); bg(ty_)
rect(ty_, 0, 6.38, 13.33, 1.12, fill=BG_CARD)
rect(ty_, 0, 6.36, 13.33, 0.06, fill=ACCENT)
txt(ty_, "🎵", 9.5, 0.8, 3.5, 3.5, size=110, align=PP_ALIGN.CENTER)
txt(ty_, "Thank You", 0.6, 1.1, 8.5, 1.2, size=64, bold=True, colour=WHITE)
txt(ty_, "MyTunes – Ad-Free Social Music Streaming Platform",
    0.6, 2.28, 8.8, 0.55, size=17, colour=ACCENT)
rect(ty_, 0.6, 3.0, 4.5, 0.05, fill=ACCENT2)

achievements = [
    "Full-stack web app: Node.js + Express + MongoDB + HTML/CSS/JS",
    "Real-time sync via Socket.io for group listening sessions",
    "Secure JWT auth, bcrypt hashing, refresh tokens implemented",
    "Mood-based discovery, friends system & social features delivered",
    "Automated test suite written with Jest + Supertest",
]
for i, line in enumerate(achievements):
    rect(ty_, 0.6, 3.22+i*0.52, 0.22, 0.22, fill=ACCENT)
    txt(ty_, line, 0.95, 3.2+i*0.52, 7.5, 0.32, size=12, colour=GREY)

txt(ty_, "Department of Computer Science  |  Academic Year 2024-25",
    0.6, 6.55, 12, 0.28, size=10, colour=GREY)


# =============================================================================
# Save
# =============================================================================
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
